#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Скрипт для извлечения адресов электронной почты из офисных документов и PDF файлов.
"""

import os
import re
import csv
from pathlib import Path
from tkinter import filedialog, messagebox, Tk
from tkinter.ttk import Progressbar
import tkinter as tk
from threading import Thread
import time

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import xlrd
    XLS_AVAILABLE = True
except ImportError:
    XLS_AVAILABLE = False


class EmailExtractor:
    """Класс для извлечения email адресов из документов."""
    
    # Регулярное выражение для поиска email адресов
    EMAIL_PATTERN = re.compile(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b')
    
    # Поддерживаемые расширения файлов
    SUPPORTED_EXTENSIONS = {
        '.docx': 'Word',
        '.doc': 'Word (старый формат)',
        '.xlsx': 'Excel',
        '.xls': 'Excel (старый формат)',
        '.pptx': 'PowerPoint',
        '.ppt': 'PowerPoint (старый формат)',
        '.pdf': 'PDF',
        '.txt': 'Текстовый файл'
    }
    
    def __init__(self, progress_callback=None, status_callback=None):
        """
        Инициализация экстрактора.
        
        Args:
            progress_callback: Функция для обновления прогресса (current, total)
            status_callback: Функция для обновления статуса (message)
        """
        self.progress_callback = progress_callback
        self.status_callback = status_callback
        self.found_emails = {}  # {email: [список файлов]}
        self.processed_files = 0
        self.total_files = 0
        self.errors = []
    
    def update_progress(self, current, total):
        """Обновить прогресс обработки."""
        if self.progress_callback:
            self.progress_callback(current, total)
    
    def update_status(self, message):
        """Обновить статус обработки."""
        if self.status_callback:
            self.status_callback(message)
    
    def find_files(self, directory):
        """
        Найти все поддерживаемые файлы в директории рекурсивно.
        
        Args:
            directory: Путь к директории для поиска
            
        Returns:
            Список путей к файлам
        """
        files = []
        directory = Path(directory)
        
        for ext in self.SUPPORTED_EXTENSIONS.keys():
            files.extend(directory.rglob(f'*{ext}'))
            files.extend(directory.rglob(f'*{ext.upper()}'))
        
        return sorted(files)
    
    def extract_emails_from_text(self, text):
        """
        Извлечь email адреса из текста.
        
        Args:
            text: Текст для поиска
            
        Returns:
            Множество найденных email адресов
        """
        if not text:
            return set()
        return set(self.EMAIL_PATTERN.findall(text))
    
    def extract_from_docx(self, file_path):
        """Извлечь email из .docx файла."""
        try:
            doc = Document(file_path)
            text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            # Также проверяем таблицы
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += '\n' + cell.text
            return self.extract_emails_from_text(text)
        except Exception as e:
            self.errors.append(f"Ошибка при обработке {file_path}: {str(e)}")
            return set()
    
    def extract_from_xlsx(self, file_path):
        """Извлечь email из .xlsx файла."""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ''
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell:
                            text += str(cell) + '\n'
            return self.extract_emails_from_text(text)
        except Exception as e:
            self.errors.append(f"Ошибка при обработке {file_path}: {str(e)}")
            return set()
    
    def extract_from_xls(self, file_path):
        """Извлечь email из .xls файла (старый формат)."""
        try:
            workbook = xlrd.open_workbook(file_path)
            text = ''
            for sheet in workbook.sheets():
                for row_idx in range(sheet.nrows):
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        if cell.value:
                            text += str(cell.value) + '\n'
            return self.extract_emails_from_text(text)
        except Exception as e:
            self.errors.append(f"Ошибка при обработке {file_path}: {str(e)}")
            return set()
    
    def extract_from_pptx(self, file_path):
        """Извлечь email из .pptx файла."""
        try:
            prs = Presentation(file_path)
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
                    # Проверяем таблицы
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text += cell.text + '\n'
            return self.extract_emails_from_text(text)
        except Exception as e:
            self.errors.append(f"Ошибка при обработке {file_path}: {str(e)}")
            return set()
    
    def extract_from_pdf(self, file_path):
        """Извлечь email из PDF файла."""
        try:
            text = ''
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + '\n'
            return self.extract_emails_from_text(text)
        except Exception as e:
            self.errors.append(f"Ошибка при обработке {file_path}: {str(e)}")
            return set()
    
    def extract_from_txt(self, file_path):
        """Извлечь email из текстового файла."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
            return self.extract_emails_from_text(text)
        except Exception as e:
            self.errors.append(f"Ошибка при обработке {file_path}: {str(e)}")
            return set()
    
    def extract_from_doc(self, file_path):
        """Извлечь email из .doc файла (старый формат)."""
        # Для старых .doc файлов нужны дополнительные библиотеки
        # Пытаемся прочитать как текстовый файл (может не сработать)
        try:
            return self.extract_from_txt(file_path)
        except:
            self.errors.append(f"Не удалось обработать {file_path}: требуется дополнительная библиотека для .doc файлов")
            return set()
    
    def extract_from_ppt(self, file_path):
        """Извлечь email из .ppt файла (старый формат)."""
        # Для старых .ppt файлов нужны дополнительные библиотеки
        self.errors.append(f"Не удалось обработать {file_path}: требуется дополнительная библиотека для .ppt файлов")
        return set()
    
    def process_file(self, file_path):
        """
        Обработать один файл и извлечь email адреса.
        
        Args:
            file_path: Путь к файлу
        """
        file_path = Path(file_path)
        ext = file_path.suffix.lower()
        
        emails = set()
        
        if ext == '.docx' and DOCX_AVAILABLE:
            emails = self.extract_from_docx(file_path)
        elif ext == '.doc':
            emails = self.extract_from_doc(file_path)
        elif ext == '.xlsx' and XLSX_AVAILABLE:
            emails = self.extract_from_xlsx(file_path)
        elif ext == '.xls' and XLS_AVAILABLE:
            emails = self.extract_from_xls(file_path)
        elif ext == '.pptx' and PPTX_AVAILABLE:
            emails = self.extract_from_pptx(file_path)
        elif ext == '.ppt':
            emails = self.extract_from_ppt(file_path)
        elif ext == '.pdf' and PDF_AVAILABLE:
            emails = self.extract_from_pdf(file_path)
        elif ext == '.txt':
            emails = self.extract_from_txt(file_path)
        else:
            self.errors.append(f"Неподдерживаемый формат или отсутствует библиотека: {file_path}")
            return
        
        # Сохраняем найденные email адреса
        for email in emails:
            email_lower = email.lower()
            if email_lower not in self.found_emails:
                self.found_emails[email_lower] = []
            self.found_emails[email_lower].append(str(file_path))
        
        self.processed_files += 1
        self.update_progress(self.processed_files, self.total_files)
    
    def process_directory(self, directory):
        """
        Обработать все файлы в директории.
        
        Args:
            directory: Путь к директории
        """
        self.update_status("Поиск файлов...")
        files = self.find_files(directory)
        self.total_files = len(files)
        
        if self.total_files == 0:
            self.update_status("Файлы не найдены!")
            return
        
        self.update_status(f"Найдено файлов: {self.total_files}. Начинаю обработку...")
        self.processed_files = 0
        self.found_emails = {}
        self.errors = []
        
        for i, file_path in enumerate(files, 1):
            self.update_status(f"Обработка: {file_path.name} ({i}/{self.total_files})")
            self.process_file(file_path)
        
        self.update_status(f"Обработка завершена! Найдено уникальных email: {len(self.found_emails)}")
    
    def save_to_csv(self, output_file):
        """
        Сохранить результаты в CSV файл.
        
        Args:
            output_file: Путь к выходному файлу
        """
        with open(output_file, 'w', newline='', encoding='utf-8-sig') as csvfile:
            writer = csv.writer(csvfile)
            writer.writerow(['Email адрес', 'Количество файлов', 'Файлы'])
            
            for email, files in sorted(self.found_emails.items()):
                files_str = '; '.join(files)
                writer.writerow([email, len(files), files_str])
    
    def save_to_excel(self, output_file):
        """
        Сохранить результаты в Excel файл.
        
        Args:
            output_file: Путь к выходному файлу
        """
        if not XLSX_AVAILABLE:
            raise ImportError("openpyxl не установлен. Используйте CSV формат.")
        
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Найденные email"
        
        # Заголовки
        ws.append(['Email адрес', 'Количество файлов', 'Файлы'])
        
        # Данные
        for email, files in sorted(self.found_emails.items()):
            files_str = '; '.join(files)
            ws.append([email, len(files), files_str])
        
        # Автоматическая ширина столбцов
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 100)
            ws.column_dimensions[column_letter].width = adjusted_width
        
        wb.save(output_file)


class EmailExtractorGUI:
    """GUI приложение для извлечения email адресов."""
    
    def __init__(self, root):
        self.root = root
        self.root.title("Извлечение email адресов из документов")
        self.root.geometry("700x500")
        
        self.selected_directory = None
        self.extractor = None
        self.is_processing = False
        
        self.setup_ui()
    
    def setup_ui(self):
        """Настройка интерфейса."""
        # Заголовок
        title_label = tk.Label(
            self.root,
            text="Извлечение email адресов из документов",
            font=("Arial", 14, "bold"),
            pady=10
        )
        title_label.pack()
        
        # Кнопка выбора папки
        self.select_button = tk.Button(
            self.root,
            text="Выбрать папку",
            command=self.select_directory,
            font=("Arial", 12),
            padx=20,
            pady=10
        )
        self.select_button.pack(pady=10)
        
        # Информация о выбранной папке
        self.directory_label = tk.Label(
            self.root,
            text="Папка не выбрана",
            font=("Arial", 10),
            wraplength=650,
            justify="left"
        )
        self.directory_label.pack(pady=5)
        
        # Кнопка запуска
        self.start_button = tk.Button(
            self.root,
            text="Начать обработку",
            command=self.start_processing,
            font=("Arial", 12),
            padx=20,
            pady=10,
            state="disabled"
        )
        self.start_button.pack(pady=10)
        
        # Прогресс-бар
        self.progress_label = tk.Label(
            self.root,
            text="",
            font=("Arial", 10)
        )
        self.progress_label.pack(pady=5)
        
        self.progress_bar = Progressbar(
            self.root,
            length=650,
            mode='determinate'
        )
        self.progress_bar.pack(pady=5)
        
        # Статус
        self.status_label = tk.Label(
            self.root,
            text="Готов к работе",
            font=("Arial", 10),
            wraplength=650,
            justify="left",
            anchor="w"
        )
        self.status_label.pack(pady=5, padx=20, fill="x")
        
        # Область для вывода результатов
        results_frame = tk.Frame(self.root)
        results_frame.pack(pady=10, padx=20, fill="both", expand=True)
        
        results_label = tk.Label(
            results_frame,
            text="Результаты:",
            font=("Arial", 10, "bold")
        )
        results_label.pack(anchor="w")
        
        self.results_text = tk.Text(
            results_frame,
            height=8,
            font=("Courier", 9),
            wrap="word"
        )
        self.results_text.pack(fill="both", expand=True)
        
        scrollbar = tk.Scrollbar(self.results_text)
        scrollbar.pack(side="right", fill="y")
        self.results_text.config(yscrollcommand=scrollbar.set)
        scrollbar.config(command=self.results_text.yview)
    
    def select_directory(self):
        """Выбрать директорию для обработки."""
        directory = filedialog.askdirectory(title="Выберите папку для поиска email адресов")
        if directory:
            self.selected_directory = directory
            self.directory_label.config(text=f"Выбрана папка: {directory}")
            self.start_button.config(state="normal")
    
    def update_progress(self, current, total):
        """Обновить прогресс-бар."""
        if total > 0:
            progress = (current / total) * 100
            self.progress_bar['value'] = progress
            self.progress_label.config(text=f"Обработано: {current} из {total} файлов ({progress:.1f}%)")
    
    def update_status(self, message):
        """Обновить статус."""
        self.status_label.config(text=message)
        self.root.update_idletasks()
    
    def log_result(self, message):
        """Добавить сообщение в область результатов."""
        self.results_text.insert("end", message + "\n")
        self.results_text.see("end")
        self.root.update_idletasks()
    
    def start_processing(self):
        """Начать обработку файлов."""
        if not self.selected_directory or self.is_processing:
            return
        
        self.is_processing = True
        self.select_button.config(state="disabled")
        self.start_button.config(state="disabled")
        self.progress_bar['value'] = 0
        self.results_text.delete(1.0, "end")
        
        # Запускаем обработку в отдельном потоке
        thread = Thread(target=self.process_files)
        thread.daemon = True
        thread.start()
    
    def process_files(self):
        """Обработать файлы (выполняется в отдельном потоке)."""
        try:
            self.extractor = EmailExtractor(
                progress_callback=self.update_progress,
                status_callback=self.update_status
            )
            
            self.log_result("=" * 60)
            self.log_result("Начало обработки...")
            self.log_result(f"Директория: {self.selected_directory}")
            self.log_result("=" * 60)
            
            self.extractor.process_directory(self.selected_directory)
            
            # Выводим результаты
            self.log_result("\n" + "=" * 60)
            self.log_result("РЕЗУЛЬТАТЫ ОБРАБОТКИ")
            self.log_result("=" * 60)
            self.log_result(f"Всего обработано файлов: {self.extractor.processed_files}")
            self.log_result(f"Найдено уникальных email адресов: {len(self.extractor.found_emails)}")
            
            if self.extractor.errors:
                self.log_result(f"\nОшибок при обработке: {len(self.extractor.errors)}")
                for error in self.extractor.errors[:10]:  # Показываем первые 10 ошибок
                    self.log_result(f"  - {error}")
                if len(self.extractor.errors) > 10:
                    self.log_result(f"  ... и еще {len(self.extractor.errors) - 10} ошибок")
            
            # Сохраняем результаты
            if self.extractor.found_emails:
                output_dir = Path(self.selected_directory)
                csv_file = output_dir / "найденные_email.csv"
                excel_file = output_dir / "найденные_email.xlsx"
                
                try:
                    self.extractor.save_to_csv(csv_file)
                    self.log_result(f"\nРезультаты сохранены в CSV: {csv_file}")
                except Exception as e:
                    self.log_result(f"\nОшибка при сохранении CSV: {str(e)}")
                
                try:
                    if XLSX_AVAILABLE:
                        self.extractor.save_to_excel(excel_file)
                        self.log_result(f"Результаты сохранены в Excel: {excel_file}")
                except Exception as e:
                    self.log_result(f"Ошибка при сохранении Excel: {str(e)}")
                
                # Показываем список найденных email
                self.log_result("\n" + "=" * 60)
                self.log_result("НАЙДЕННЫЕ EMAIL АДРЕСА:")
                self.log_result("=" * 60)
                for email, files in sorted(self.extractor.found_emails.items()):
                    self.log_result(f"{email} (найден в {len(files)} файле(ах))")
            else:
                self.log_result("\nEmail адреса не найдены.")
            
            self.log_result("\n" + "=" * 60)
            self.log_result("Обработка завершена!")
            self.log_result("=" * 60)
            
            messagebox.showinfo(
                "Завершено",
                f"Обработка завершена!\n\n"
                f"Обработано файлов: {self.extractor.processed_files}\n"
                f"Найдено email: {len(self.extractor.found_emails)}\n"
                f"Результаты сохранены в выбранной папке."
            )
        except Exception as e:
            self.log_result(f"\nКРИТИЧЕСКАЯ ОШИБКА: {str(e)}")
            messagebox.showerror("Ошибка", f"Произошла ошибка: {str(e)}")
        finally:
            self.is_processing = False
            self.select_button.config(state="normal")
            self.start_button.config(state="normal")


def main():
    """Главная функция."""
    # Проверяем доступность библиотек
    missing_libs = []
    if not DOCX_AVAILABLE:
        missing_libs.append("python-docx")
    if not XLSX_AVAILABLE:
        missing_libs.append("openpyxl")
    if not PDF_AVAILABLE:
        missing_libs.append("PyPDF2")
    if not PPTX_AVAILABLE:
        missing_libs.append("python-pptx")
    if not XLS_AVAILABLE:
        missing_libs.append("xlrd")
    
    root = Tk()
    app = EmailExtractorGUI(root)
    
    if missing_libs:
        messagebox.showwarning(
            "Предупреждение",
            f"Некоторые библиотеки не установлены:\n{', '.join(missing_libs)}\n\n"
            f"Установите их командой:\npip install {' '.join(missing_libs)}\n\n"
            f"Некоторые форматы файлов могут не обрабатываться."
        )
    
    root.mainloop()


if __name__ == "__main__":
    main()

